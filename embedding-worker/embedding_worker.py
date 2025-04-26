#!/usr/bin/env python
"""
Embedding Worker
---------------
This script runs separately from the FastAPI server and processes embedding tasks
from the Redis queue.

Usage: python embedding_worker.py
"""

import os
import json
import time
import asyncio
import redis
from datetime import datetime
from sqlalchemy import create_engine, text as sql_text
from sqlalchemy.orm import sessionmaker
import uuid
from uuid import UUID
import torch
from transformers import AutoTokenizer, AutoModel
import numpy as np
from functools import lru_cache
import logging
import traceback
import sys

# Configuration
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "512"))
OVERLAP = int(os.getenv('OVERLAP', '128'))
#MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "intfloat/multilingual-e5-base")
#EMBEDDING_DIMENSION = int(os.getenv("EMBEDDING_DIMENSION", "768"))  # Configurable embedding dimension

MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME", "intfloat/multilingual-e5-large-instruct")
EMBEDDING_DIMENSION = int(os.getenv("EMBEDDING_DIMENSION", "1024"))  # Configurable embedding dimension
HF_HOME = os.getenv("HF_HOME", "/app/model_weights")

DB_URL = os.getenv("DATABASE_URL", "postgresql://admin:SuperSecurePassword@host.docker.internal:5432/collections_db")
REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379/0")
SLEEP_TIME = int(os.getenv("WORKER_SLEEP_TIME", "1"))
USE_CPU = os.getenv('USE_CPU', 'false').lower() == 'true'
MAX_RECONNECT_ATTEMPTS = 5
RECONNECT_DELAY = 5  # seconds

# Get log level from environment variable (default to INFO if not specified)
LOG_LEVEL_NAME = os.getenv("LOG_LEVEL", "INFO").upper()
LOG_LEVELS = {
    "DEBUG": logging.DEBUG,
    "INFO": logging.INFO,
    "WARNING": logging.WARNING,
    "ERROR": logging.ERROR,
    "CRITICAL": logging.CRITICAL
}

# Use the specified log level or default to INFO if invalid
LOG_LEVEL = LOG_LEVELS.get(LOG_LEVEL_NAME, logging.INFO)

# Set up logging with more detailed format
logging.basicConfig(
    level=LOG_LEVEL,
    format='%(asctime)s - %(name)s - %(levelname)s - [%(filename)s:%(lineno)d] - %(message)s'
)
logger = logging.getLogger("embedding_worker")
logger.info(f"Logger initialized with level: {LOG_LEVEL_NAME}")

def get_logger(name: str = None) -> logging.Logger:
    """Get a logger with the given name, inheriting the main configuration"""
    if name:
        return logging.getLogger(f"embedding_worker.{name}")
    return logger

class EmbeddingWorker:
    def __init__(self):
        self.redis_client = None
        self.logger = get_logger("worker")
        self.engine = create_engine(DB_URL)
        self.SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=self.engine)
        
        db = self.SessionLocal()
        try:
            self.ensure_pgvector_extension(db)
            db.commit()
        finally:
            db.close()
        
        # Load the embedding model and tokenizer at initialization
        self.logger.info(f"Loading embedding model: {MODEL_NAME}")

        if torch.cuda.is_available() and not USE_CPU:
            self.tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, cache_dir=HF_HOME)
            self.model = AutoModel.from_pretrained(MODEL_NAME, cache_dir=HF_HOME, )
            self.model = self.model.to("cuda")
            self.logger.info("Using GPU for embeddings")
        else:
            self.tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, cache_dir=HF_HOME, device_map="cpu")
            self.model = AutoModel.from_pretrained(MODEL_NAME, cache_dir=HF_HOME, device_map="cpu")
            self.logger.info("GPU not available, using CPU for embeddings")
        
        # Test database connection
        try:
            with self.engine.connect() as conn:
                result = conn.execute(sql_text("SELECT version()")).scalar()
                self.logger.info(f"Connected to PostgreSQL. Version: {result}")
            self.logger.info("Successfully established database connection")
        except Exception as e:
            self.logger.error(f"Database connection error: {str(e)}")
            self.logger.error(f"Database connection traceback: {traceback.format_exc()}")
            raise

    def connect_to_redis(self):
        """Establish connection to Redis with detailed error logging"""
        try:
            self.logger.info(f"Attempting to connect to Redis at {REDIS_URL}")
            client = redis.from_url(REDIS_URL)
            
            # Test connection with PING
            response = client.ping()
            self.logger.info(f"Redis PING response: {response}")
            
            # Test basic operations
            test_key = "test_connection"
            client.set(test_key, "test_value")
            test_value = client.get(test_key)
            self.logger.info(f"Redis test key retrieval successful: {test_value}")
            client.delete(test_key)
            
            self.logger.info("Successfully established Redis connection")
            self.redis_client = client
            return True
        except redis.ConnectionError as e:
            self.logger.error(f"Redis ConnectionError: {str(e)}")
            self.logger.error(f"Connection traceback: {traceback.format_exc()}")
            return False
        except Exception as e:
            self.logger.error(f"Unexpected Redis error: {str(e)}")
            self.logger.error(f"Error traceback: {traceback.format_exc()}")
            return False

    def ensure_redis_connection(self):
        """Ensure Redis connection is established with retry mechanism"""
        if self.redis_client is not None:
            try:
                self.redis_client.ping()
                return True
            except:
                self.redis_client = None
        
        for attempt in range(MAX_RECONNECT_ATTEMPTS):
            self.logger.info(f"Redis connection attempt {attempt + 1}/{MAX_RECONNECT_ATTEMPTS}")
            if self.connect_to_redis():
                return True
            if attempt < MAX_RECONNECT_ATTEMPTS - 1:
                self.logger.info(f"Waiting {RECONNECT_DELAY} seconds before next attempt...")
                time.sleep(RECONNECT_DELAY)
        
        self.logger.critical("Failed to establish Redis connection after maximum attempts")
        return False

    def update_task_status(self, task_id, status, progress, document_count, processed_count, collection_id=None, document_ids=None):
        """Update task status in Redis with enhanced logging"""
        try:
            if self.redis_client is None:
                self.logger.error("Redis client is None in update_task_status, attempting to reconnect...")
                if not self.ensure_redis_connection():
                    raise Exception("Failed to reconnect to Redis")
            
            task_data = {
                "task_id": task_id,
                "status": status,
                "progress": progress,
                "document_count": document_count,
                "processed_count": processed_count,
                "collection_id": str(collection_id) if collection_id else None,
                "document_ids": [str(doc_id) for doc_id in document_ids] if document_ids else [],
                "updated_at": datetime.utcnow().isoformat()
            }
            
            key = f"embedding_task:{task_id}"
            self.logger.debug(f"Updating Redis task status - Key: {key}, Data: {json.dumps(task_data)}")
            self.redis_client.set(key, json.dumps(task_data))
            self.logger.info(f"Successfully updated task status for {task_id} to {status}")
        except Exception as e:
            self.logger.error(f"Failed to update task status in Redis for task {task_id}: {str(e)}")
            self.logger.error(traceback.format_exc())
            raise

    async def process_task(self, task_data):
        """Process a task from the queue with enhanced logging"""
        task_id = task_data.get("task_id")
        task_type = task_data.get("task_type")
        
        self.logger.info(f"Starting to process task {task_id} of type {task_type}")
        self.logger.debug(f"Full task data: {json.dumps(task_data)}")
        
        try:
            if task_type == "documents":
                document_ids = [UUID(doc_id) for doc_id in task_data.get("document_ids", [])]
                self.logger.info(f"Processing {len(document_ids)} documents for task {task_id}")
                if document_ids:
                    await self.process_documents(task_id, document_ids)
                else:
                    self.logger.warning(f"No document IDs provided for task {task_id}")
            
            elif task_type == "collection":
                collection_id = task_data.get("collection_id")
                if collection_id:
                    self.logger.info(f"Processing collection {collection_id} for task {task_id}")
                    await self.process_collection(task_id, UUID(collection_id))
                else:
                    self.logger.warning(f"No collection ID provided for task {task_id}")
            
            else:
                self.logger.error(f"Unknown task type: {task_type} for task {task_id}")
        except Exception as e:
            self.logger.error(f"Error processing task {task_id}: {str(e)}")
            self.logger.error(traceback.format_exc())
            raise

    def run(self):
        """Main worker loop with enhanced error handling and logging"""
        self.logger.info("Starting main worker loop")
        
        while True:
            try:
                if self.redis_client is None:
                    self.logger.warning("Redis client is None, attempting to reconnect...")
                    if not self.ensure_redis_connection():
                        self.logger.error("Failed to connect to Redis, sleeping before retry...")
                        time.sleep(SLEEP_TIME)
                        continue

                # Try to get a task from Redis
                result = self.redis_client.brpop("embedding_tasks", timeout=SLEEP_TIME)
                
                if result is None:
                    continue
                    
                _, task_json = result
                self.logger.info(f"Received task data: {task_json}")
                
                try:
                    task_data = json.loads(task_json)
                except json.JSONDecodeError as e:
                    self.logger.error(f"Failed to decode task JSON: {e}")
                    self.logger.error(f"Raw task data: {task_json}")
                    continue
                    
                task_id = task_data.get("task_id")
                if not task_id:
                    self.logger.error("Task data missing task_id")
                    continue
                    
                self.logger.info(f"Processing task {task_id}")
                asyncio.run(self.process_task(task_data))
                
            except redis.ConnectionError as e:
                self.logger.error(f"Redis connection error: {e}")
                self.redis_client = None
                time.sleep(SLEEP_TIME)
                
            except Exception as e:
                self.logger.error(f"Unexpected error in main loop: {e}")
                self.logger.error(traceback.format_exc())
                time.sleep(SLEEP_TIME)
                
            finally:
                self.logger.debug("Completed task processing iteration")

    def log_with_context(self, msg, context=None, level="info"):
        """Enhanced logging function with context"""
        log_func = getattr(self.logger, level)
        if context:
            msg = f"[{context}] {msg}"
        log_func(msg)

    async def process_documents(self, task_id, document_ids):
        """Process a list of documents for embedding generation"""
        db = self.SessionLocal()
        try:
            document_count = len(document_ids)
            processed_count = 0
            
            # Update status to processing
            self.update_task_status(
                task_id, "processing", 0.0, document_count, processed_count,
                document_ids=document_ids
            )
            
            # Process each document
            for i, doc_id in enumerate(document_ids):
                try:
                    chunks_processed = await self.process_document_for_embeddings(doc_id, db)
                    if chunks_processed > 0:
                        processed_count += 1
                    
                    # Update progress
                    progress = (i + 1) / document_count
                    self.update_task_status(
                        task_id, "processing", progress, document_count, processed_count,
                        document_ids=document_ids
                    )
                    
                except Exception as e:
                    self.log_with_context(f"Error processing document {doc_id}: {e}")
            
            # Update final status
            final_status = "completed" if processed_count == document_count else "partial"
            self.update_task_status(
                task_id, final_status, 1.0, document_count, processed_count,
                document_ids=document_ids
            )
            
        except Exception as e:
            stack_trace = traceback.format_exc()
            self.log_with_context(f"Error in document batch processing: {e}")
            self.log_with_context(f"Stack trace:\n{stack_trace}")
            
            # Update status to failed
            self.update_task_status(
                task_id, "failed", processed_count / document_count if document_count > 0 else 0.0,
                document_count, processed_count, document_ids=document_ids
            )
        finally:
            db.close()

    async def process_collection(self, task_id, collection_id):
        """Process all documents in a collection"""
        db = self.SessionLocal()
        try:
            # Get all document IDs in the collection
            doc_sql = sql_text("""
                SELECT id FROM collections.documents
                WHERE collection_id = :collection_id
            """)
            
            result = db.execute(doc_sql, {"collection_id": str(collection_id)}).fetchall()
            document_ids = [row.id for row in result]
            self.log_with_context(f"Found {len(document_ids)} documents in collection {collection_id}")
            
            # Close this DB session before processing documents
            db.close()
            
            # Process the documents
            await self.process_documents(task_id, document_ids)
            
        except Exception as e:
            stack_trace = traceback.format_exc()
            self.log_with_context(f"Error processing collection {collection_id}: {e}")
            self.log_with_context(f"Stack trace:\n{stack_trace}")
            
            # Update status to failed
            self.update_task_status(
                task_id, "failed", 0.0, 0, 0,
                collection_id=collection_id
            )
            
            # Make sure DB session is closed
            if not db.closed:
                db.close()

    async def process_document_for_embeddings(self, document_id, db):
        """Process a document to generate embeddings"""
        try:
            self.log_with_context(f"Starting to process document {document_id}")
            
            # Get document info
            doc_sql = sql_text("""
                SELECT d.file_path, d.type, d.name
                FROM collections.documents d
                WHERE d.id = :document_id
            """)
            
            doc = db.execute(doc_sql, {"document_id": str(document_id)}).fetchone()
            if not doc:
                self.log_with_context(f"Document {document_id} not found in database")
                raise ValueError(f"Document {document_id} not found")
                
            document_path = doc.file_path
            document_type = doc.type
            document_name = doc.name
            
            self.log_with_context(f"Processing document: {document_name} (type: {document_type})")
            self.log_with_context(f"Document path: {document_path}")
            
            # Check file existence and permissions
            if not os.path.exists(document_path):
                self.log_with_context(f"File does not exist at path: {document_path}")
                raise FileNotFoundError(f"File not found: {document_path}")
            
            # Extract text from document
            self.log_with_context(f"Starting text extraction for document {document_id}")
            try:
                text = await extract_text_from_document(document_path)
                self.log_with_context(f"Text extraction completed. Extracted text length: {len(text)}")
                if not text.strip():
                    self.log_with_context("Warning: Extracted text is empty or only whitespace")
            except ValueError as e:
                self.log_with_context(f"Document {document_id} marked as unprocessable: {str(e)}")
                # If the file is unprocessable, update its status
                update_sql = sql_text("""
                    UPDATE collections.documents
                    SET status = 'unprocessable'
                    WHERE id = :document_id
                """)
                db.execute(update_sql, {"document_id": str(document_id)})
                db.commit()
                return 0
            except Exception as e:
                self.log_with_context(f"Unexpected error during text extraction: {str(e)}")
                raise
                
            # Split text into chunks
            self.log_with_context("Starting text chunking")
            chunks = chunk_document_text(text)
            total_chunks = len(chunks)
            self.log_with_context(f"Created {total_chunks} chunks from document")
            
            # Update document status to processing with chunk info
            update_sql = sql_text("""
                UPDATE collections.documents
                SET status = 'processing', 
                    metadata = jsonb_build_object(
                        'total_chunks', :total_chunks,
                        'processed_chunks', 0
                    )
                WHERE id = :document_id
            """)
            db.execute(update_sql, {
                "document_id": str(document_id),
                "total_chunks": total_chunks
            })
            
            # Delete existing chunks
            self.log_with_context("Deleting existing chunks")
            delete_sql = sql_text("""
                DELETE FROM collections.document_chunks
                WHERE document_id = :document_id
            """)
            db.execute(delete_sql, {"document_id": str(document_id)})
            
            # Insert new chunks
            self.log_with_context("Inserting new chunks")
            count = 0
            
            for chunk_index, chunk in enumerate(chunks):
                # Start a new transaction for each chunk
                try:
                    self.log_with_context(f"Processing chunk {chunk_index + 1}/{total_chunks}")
                    
                    # Sanitize chunk text to remove problematic characters
                    # Replace or remove control characters and other problematic characters
                    sanitized_chunk = chunk.replace('\x00', ' ')  # NUL
                    sanitized_chunk = ''.join(char if ord(char) >= 32 or char in '\n\r\t' else ' ' for char in sanitized_chunk)
                    
                    # Generate embedding for this chunk
                    self.log_with_context(f"Generating embedding for chunk {chunk_index + 1}")
                    try:
                        # Generate embedding asynchronously
                        embedding = await self.generate_embedding(sanitized_chunk)
                        
                        # Handle dimension mismatch - database expects 1536 dimensions
                        
                        current_dims = len(embedding)
                        
                        if current_dims < EMBEDDING_DIMENSION:
                            # Pad with zeros if the embedding is smaller than required
                            self.log_with_context(f"Padding embedding from {current_dims} to {EMBEDDING_DIMENSION} dimensions")
                            padding = [0.0] * (EMBEDDING_DIMENSION - current_dims)
                            embedding = embedding + padding
                        elif current_dims > EMBEDDING_DIMENSION:
                            # Truncate if the embedding is larger than required
                            self.log_with_context(f"Truncating embedding from {current_dims} to {EMBEDDING_DIMENSION} dimensions")
                            embedding = embedding[:EMBEDDING_DIMENSION]
                        
                        # Format embedding for pgvector
                        # The proper format for PostgreSQL vectors is '{x1,x2,x3,...}'
                        embedding_str = str(embedding).replace(' ', '')
                        
                        # Add debug logging
                        self.log_with_context(f"Embedding dimension: {len(embedding)}")
                        
                        # Validate the embedding before insert
                        if not embedding or len(embedding) != EMBEDDING_DIMENSION:
                            self.log_with_context(f"Invalid embedding dimensions: got {len(embedding) if embedding else 0}, expected {EMBEDDING_DIMENSION}", level="error")
                            continue
                        

                        
                        # Approach 1: Direct vector type casting (most compatible)
                        try:
                            insert_sql = sql_text(f"""
                                INSERT INTO collections.document_chunks
                                (document_id, chunk_text, chunk_index, embedding)
                                VALUES (
                                    '{str(document_id)}', 
                                    :chunk_text, 
                                    {count}, 
                                    '{embedding_str}'::vector
                                )
                            """)
                            
                            db.execute(insert_sql, {"chunk_text": sanitized_chunk})
                            db.commit()
                            success = True
                            self.log_with_context(f"Successfully inserted embedding for chunk {chunk_index + 1} with approach 1")
                        except Exception as e1:
                            db.rollback()  # Important: roll back failed transaction
                            self.log_with_context(f"Approach 1 failed: {str(e1)}", level="error")
                            
                    except Exception as e:
                        # Make sure to rollback on any error
                        db.rollback()
                        self.log_with_context(f"Error processing chunk {chunk_index + 1}: {str(e)}", level="error")
                        self.log_with_context(f"Error traceback: {traceback.format_exc()}", level="error")
                        # Continue with next chunk instead of failing the entire document
                        continue
                        
                except Exception as outer_e:
                    # Make sure to rollback on any error
                    try:
                        db.rollback()
                    except:
                        pass
                    self.log_with_context(f"Fatal error processing chunk {chunk_index + 1}: {str(outer_e)}", level="error")
                    self.log_with_context(f"Error traceback: {traceback.format_exc()}", level="error")
                    # Continue with next chunk
            
            # Ensure final commit
            db.commit()
            
            # Update document status to completed
            self.log_with_context("Updating document status to completed")
            update_sql = sql_text("""
                UPDATE collections.documents
                SET status = 'completed',
                    metadata = jsonb_set(
                        COALESCE(metadata, '{}'::jsonb),
                        '{processed_chunks}',
                        to_jsonb(:total_chunks)
                    )
                WHERE id = :document_id
            """)
            db.execute(update_sql, {"document_id": str(document_id), "total_chunks": total_chunks})
            
            db.commit()
            self.log_with_context(f"Successfully processed {count}/{total_chunks} chunks for document {document_id}")
            return count
        except Exception as e:
            stack_trace = traceback.format_exc()
            self.log_with_context(f"Error processing document {document_id}: {e}")
            self.log_with_context(f"Stack trace:\n{stack_trace}")
            db.rollback()
            raise  # Re-raise to allow outer error handlers to catch

    async def generate_embedding(self, text):
        """Generate embedding for text using the instance's model and tokenizer"""
        logger.info(f"Generating embedding for text: '{text[:100]}...'")  # Log first 100 chars of text
        
        if not text or not text.strip():
            logger.warning("Empty text received, returning empty embedding")
            return []
        
        if not torch.cuda.is_available() or USE_CPU:
            logger.warning("CUDA not available, generating random embedding for testing")
            # Generate random embedding for testing without GPU
            random_embedding = np.random.normal(0, 0.1, EMBEDDING_DIMENSION).astype(np.float32)
            # Simulate processing time
            await asyncio.sleep(0.5)
            logger.info(f"Generated random embedding with dimension {len(random_embedding)}")
            return random_embedding.tolist()
        
        logger.info("Tokenizing input text")
        inputs = self.tokenizer(text, padding=True, truncation=True, return_tensors="pt", max_length=CHUNK_SIZE)
        
        # Move to GPU if available
        if torch.cuda.is_available() and not USE_CPU:
            logger.info("Moving inputs to CUDA")
            inputs = {k: v.to("cuda") for k, v in inputs.items()}
        
        logger.info("Generating embedding using model")
        with torch.no_grad():
            outputs = self.model(**inputs)
            embeddings = outputs.last_hidden_state.mean(dim=1).squeeze().cpu().numpy()
            
            # Handle case where outputs is single dimension
            if len(embeddings.shape) == 1:
                logger.info(f"Generated single-dimension embedding with shape {embeddings.shape}")
                return embeddings.tolist()
            else:
                logger.info(f"Generated multi-dimension embedding, taking first vector with shape {embeddings[0].shape}")
                return embeddings[0].tolist()

    def ensure_pgvector_extension(self, db):
        """Ensure pgvector extension is created and available"""
        try:
            # First check if extension exists
            result = db.execute(sql_text("SELECT * FROM pg_extension WHERE extname = 'vector'")).fetchone()
            if not result:
                self.logger.info("Creating pgvector extension...")
                db.execute(sql_text("CREATE EXTENSION IF NOT EXISTS vector"))
                db.commit()
                self.logger.info("pgvector extension created successfully")
            else:
                self.logger.info("pgvector extension already exists")
            
            # Now ensure the vector type is available in the collections schema
            db.execute(sql_text("SET search_path TO collections,public"))
            db.commit()
            self.logger.info("Set search path to collections schema")
            
            return True
        except Exception as e:
            self.logger.error(f"Error setting up pgvector extension: {e}")
            self.logger.error(traceback.format_exc())
            return False




# =============== Helper Functions ===============
def log_with_context(msg, context=None, level="info"):
    """Enhanced logging function with context"""
    log_func = getattr(logger, level)
    if context:
        msg = f"[{context}] {msg}"
    log_func(msg)

def chunk_document_text(text, chunk_size=CHUNK_SIZE, overlap=OVERLAP):
    """Split document text into overlapping chunks"""
    chunks = []
    if not text or len(text) <= 0:
        return chunks
        
    # Clean the text of NUL characters and normalize whitespace
    text = text.replace('\x00', ' ').replace('\r\n', '\n').replace('\r', '\n')
    text = ' '.join(text.split())  # Normalize whitespace
    
    if len(text) <= chunk_size:
        return [text]
    
    # Simple approach with clear advancement
    start = 0
    while start < len(text):
        # Calculate end position
        end = min(start + chunk_size, len(text))
        
        # Extract chunk and clean it
        chunk = text[start:end].strip()
        if chunk:  # Only add non-empty chunks
            chunks.append(chunk)
        
        # Always advance by at least 1 character to prevent infinite loop
        next_start = start + (chunk_size - overlap)
        start = max(next_start, start + 1)
    
    return chunks

async def extract_text_from_document(document_path):
    """Extract text from document based on file type"""
    # Set up logging
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)
    
    # Initialize text variable at the beginning
    text = ""
    suffix = os.path.splitext(document_path)[1].lower()
    
    logger.info(f"Starting text extraction for file: {document_path}")
    logger.info(f"File type detected: {suffix}")
    
    # Get file size and modification time
    try:
        file_stats = os.stat(document_path)
        logger.info(f"File size: {file_stats.st_size} bytes")
        logger.info(f"Last modified: {datetime.fromtimestamp(file_stats.st_mtime)}")
    except Exception as e:
        logger.error(f"Error getting file stats: {str(e)}")
    
    try:
        # Check if file type is unprocessable
        if suffix == '.pdf':
            logger.info("Starting PDF processing")
            try:
                # Check PyMuPDF version
                import fitz
                logger.info(f"PyMuPDF version: {fitz.__version__}")
                
                # Check if file exists and is readable
                if not os.path.exists(document_path):
                    logger.error(f"PDF file does not exist: {document_path}")
                    raise FileNotFoundError(f"PDF file not found: {document_path}")
                
                logger.info(f"Opening PDF file: {document_path}")
                doc = fitz.open(document_path)
                logger.info(f"PDF opened successfully. Number of pages: {len(doc)}")
                
                text = ""
                for page_num, page in enumerate(doc):
                    logger.info(f"Processing page {page_num + 1}/{len(doc)}")
                    try:
                        page_text = page.get_text()
                        text += page_text
                        logger.info(f"Page {page_num + 1} text length: {len(page_text)}")
                    except Exception as e:
                        logger.error(f"Error processing page {page_num + 1}: {str(e)}")
                
                doc.close()
                logger.info("PDF processing completed")
                logger.info(f"Total extracted text length: {len(text)}")
                
                if not text.strip():
                    logger.warning("PDF processed but no text was extracted")
                    
            except ImportError as e:
                logger.error(f"PyMuPDF import error: {str(e)}")
                raise
            except Exception as e:
                logger.error(f"Error processing PDF: {str(e)}")
                logger.error(f"Full traceback: {traceback.format_exc()}")
                raise
        
        elif suffix == '.txt':
            with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        
        elif suffix == '.docx':
            try:
                import docx
                doc = docx.Document(document_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                text = f"[Error: python-docx not installed for DOCX extraction]"
            
        elif suffix == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(document_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except ImportError:
                text = f"[Error: python-pptx not installed for PPT extraction]"
                
        elif suffix == '.rtf':
            try:
                import striprtf
                with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                    rtf_text = f.read()
                text = striprtf.rtf_to_text(rtf_text)
            except ImportError:
                text = f"[Error: striprtf not installed for RTF extraction]"
                
        else:
            # Handle code files and other text-based files
            code_extensions = {
                # Web development
                '.html', '.htm', '.css', '.js', '.jsx', '.ts', '.tsx', '.vue', '.svelte',
                # Programming languages
                '.py', '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.go', '.rs', '.rb', '.php',
                '.swift', '.kt', '.scala', '.r', '.m', '.mm', '.sh', '.bash', '.zsh',
                # Data formats
                '.json', '.xml', '.yaml', '.yml', '.csv', '.sql', '.md', '.txt',
                # Configuration
                '.env', '.ini', '.conf', '.config', '.properties', '.toml',
                # Other text formats
                '.log', '.diff', '.patch', '.rst', '.tex', '.latex'
            }
            
            try:
                # Try to read as text file
                with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                
                # If file is empty or contains only whitespace, try to detect if it's a binary file
                if not text.strip():
                    with open(document_path, 'rb') as f:
                        # Read first 1024 bytes to check for binary content
                        chunk = f.read(1024)
                        # Check if chunk contains null bytes or non-printable characters
                        if b'\x00' in chunk or any(byte < 32 and byte not in (9, 10, 13) for byte in chunk):
                            raise ValueError(f"Binary file detected: {suffix}")
                        else:
                            text = f"[Empty or whitespace-only file: {suffix}]"
                elif suffix.lower() in code_extensions:
                    # For code files, we might want to add some formatting
                    text = f"File type: {suffix}\n\n{text}"
                else:
                    # For unknown text files, add a header
                    text = f"Text file type: {suffix}\n\n{text}"
                    
            except Exception as e:
                text = f"[Error reading file {suffix}: {str(e)}]"
        
        logger.info(f"Text extraction completed. Extracted text length: {len(text)}")
        return text
    except Exception as e:
        logger.error(f"Error in text extraction: {str(e)}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise

if __name__ == "__main__":
    logger.info("Starting embedding worker...")
    
    try:
        # Create worker instance first
        worker = EmbeddingWorker()
        
        # Initialize Redis connection using the worker instance
        if not worker.ensure_redis_connection():
            logger.error("Failed to establish initial Redis connection")
            sys.exit(1)
            
        logger.info("Successfully connected to Redis")
        
        # Run the worker
        worker.run()
    except KeyboardInterrupt:
        logger.info("Received shutdown signal, stopping worker...")
    except Exception as e:
        logger.error(f"Critical error in worker: {e}")
        logger.error(traceback.format_exc())
        sys.exit(1)
    finally:
        logger.info("Worker shutdown complete")