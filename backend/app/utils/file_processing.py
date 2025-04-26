import logging
from typing import List, Dict, Any, Optional
import io
import base64
import PyPDF2
import docx
import pptx
from striprtf.striprtf import rtf_to_text
import uuid
from datetime import datetime
from sqlalchemy.orm import Session
from database import get_chat_db
from sqlalchemy import text
# Attempt to import optional dependencies


logger = logging.getLogger(__name__)

# Common text-based file extensions
TEXT_EXTENSIONS = {
    '.txt', '.md', '.py', '.js', '.ts', '.html', '.css', '.json', '.yaml', '.yml',
    '.csv', '.java', '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.php', '.rb',
    '.swift', '.kt', '.kts', '.scala', '.rs', '.sh', '.toml', '.xml'
}

def process_file_content(file_data: Dict[str, Any], chat_id: Optional[str] = None, db: Optional[Session] = None) -> str:
    """
    Reads the content of a single file based on its extension and data format.
    Expects a dictionary with 'name' and 'content' keys.
    'content' should be a string (plain text or base64 encoded for binary).
    
    If chat_id is provided and db session is passed, the extracted text will be saved to the document store.
    """
    filename = file_data.get("name", "unknown_file")
    content_input = file_data.get("content") # Could be string (text/base64) or null
    extracted_text = ""
    logger.info(f"Processing file from payload: {filename}")

    try:
        # Determine file extension
        file_ext = '.' + filename.split('.')[-1].lower() if '.' in filename else ''

        # Handle based on content type provided or file extension
        if content_input is None:
            logger.warning(f"No content provided for file: {filename}. Skipping.")
            return "[File content not provided]"

        if not isinstance(content_input, str):
            logger.error(f"Invalid content type for file {filename}: expected string (text or base64), got {type(content_input)}. Skipping.")
            return "[Invalid file content format]"

        # --- Start Processing Logic ---
        if file_ext in TEXT_EXTENSIONS:
            logger.debug(f"Reading '{filename}' as a text file.")
            extracted_text = content_input # Assume it's already decoded text

        elif file_ext == '.pdf':
            logger.debug(f"Reading '{filename}' as a PDF file using PyPDF2.")
            try:
                # Decode base64 content to bytes
                pdf_bytes = base64.b64decode(content_input)
                pdf_file_like = io.BytesIO(pdf_bytes)
                
                # Use PyPDF2 to read the PDF
                pdf_reader = PyPDF2.PdfReader(pdf_file_like)
                text_parts = []
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    try:
                        text_parts.append(page.extract_text() or "") # Add empty string if extract_text returns None
                    except Exception as page_error: # Catch errors during text extraction per page
                        logger.warning(f"Error extracting text from page {page_num+1} of PDF '{filename}': {page_error}")
                        text_parts.append(f"[Error on page {page_num+1}]" )
                
                extracted_text = "\\n".join(text_parts)
                logger.debug(f"Successfully extracted text from PDF '{filename}' using PyPDF2. Length: {len(extracted_text)}")
            
            except base64.binascii.Error as b64_error:
                logger.error(f"Base64 decoding failed for PDF file '{filename}': {b64_error}")
                extracted_text = f"[Error processing PDF: Invalid base64 data]"
            except PyPDF2.errors.PdfReadError as pdf_read_error:
                logger.error(f"PyPDF2 could not read PDF file '{filename}': {pdf_read_error}", exc_info=True)
                extracted_text = f"[Error reading PDF: {str(pdf_read_error)}]"
            except Exception as pdf_error:
                logger.error(f"Unexpected error reading PDF file '{filename}' with PyPDF2: {pdf_error}", exc_info=True)
                extracted_text = f"[Error processing PDF file: {str(pdf_error)}]"

        elif file_ext == '.docx':
            if docx:
                logger.debug(f"Reading '{filename}' as a DOCX file using python-docx.")
                try:
                    # Decode base64 and read from bytes using io.BytesIO
                    docx_bytes = base64.b64decode(content_input)
                    document = docx.Document(io.BytesIO(docx_bytes))
                    text_parts = [p.text for p in document.paragraphs]
                    extracted_text = "\\n".join(text_parts)
                    logger.debug(f"Successfully extracted text from DOCX '{filename}'. Length: {len(extracted_text)}")
                except base64.binascii.Error as b64_error:
                    logger.error(f"Base64 decoding failed for DOCX file '{filename}': {b64_error}")
                    extracted_text = f"[Error processing DOCX: Invalid base64 data]"
                except Exception as docx_error:
                    logger.error(f"Error reading DOCX file '{filename}' with python-docx: {docx_error}", exc_info=True)
                    extracted_text = f"[Error processing DOCX file: {str(docx_error)}]"
            else:
                logger.warning(f"python-docx not installed. Cannot process DOCX file: {filename}")
                extracted_text = "[DOCX processing skipped: python-docx not installed]"

        elif file_ext == '.pptx':
            logger.debug(f"Reading '{filename}' as a PPTX file using python-pptx.")
            try:
                # Decode base64 and read from bytes using io.BytesIO
                pptx_bytes = base64.b64decode(content_input)
                presentation = pptx.Presentation(io.BytesIO(pptx_bytes))
                text_parts = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                extracted_text = "\\n".join(text_parts)
                logger.debug(f"Successfully extracted text from PPTX '{filename}'. Length: {len(extracted_text)}")
            except base64.binascii.Error as b64_error:
                logger.error(f"Base64 decoding failed for PPTX file '{filename}': {b64_error}")
                extracted_text = f"[Error processing PPTX: Invalid base64 data]"
            except Exception as pptx_error:
                logger.error(f"Error reading PPTX file '{filename}' with python-pptx: {pptx_error}", exc_info=True)
                extracted_text = f"[Error processing PPTX file: {str(pptx_error)}]"

        elif file_ext == '.rtf':
            logger.debug(f"Reading '{filename}' as an RTF file using striprtf.")
            try:
                # striprtf expects a string, so decode base64 then decode bytes to string
                rtf_bytes = base64.b64decode(content_input)
                # RTF standard encoding is often Windows-1252 or latin-1, but try UTF-8 first
                try:
                    rtf_string_content = rtf_bytes.decode('utf-8')
                except UnicodeDecodeError:
                    logger.warning(f"UTF-8 decoding failed for RTF '{filename}', trying cp1252.")
                    try:
                        rtf_string_content = rtf_bytes.decode('cp1252') # Windows default
                    except UnicodeDecodeError:
                        logger.warning(f"cp1252 decoding failed for RTF '{filename}', trying latin-1.")
                        rtf_string_content = rtf_bytes.decode('latin-1', errors='ignore') # Fallback

                # Use striprtf to convert RTF to plain text
                extracted_text = rtf_to_text(rtf_string_content)
                logger.debug(f"Successfully extracted text from RTF '{filename}'. Length: {len(extracted_text)}")
            except base64.binascii.Error as b64_error:
                logger.error(f"Base64 decoding failed for RTF file '{filename}': {b64_error}")
                extracted_text = f"[Error processing RTF: Invalid base64 data]"
            except Exception as rtf_error:
                logger.error(f"Error reading RTF file '{filename}' with striprtf: {rtf_error}", exc_info=True)
                extracted_text = f"[Error processing RTF file: {str(rtf_error)}]"

        else:
            logger.warning(f"Unsupported file type or extension '{file_ext}' for file: {filename}. Treating as plain text.")
            # Assume content_input is plain text if not handled above
            extracted_text = content_input
            if not extracted_text:
                logger.warning(f"File '{filename}' (treated as text) has empty content.")
                
        # If chat_id and db session are provided and we have extracted text, save to document store

        if chat_id and db and extracted_text:
            try:
                # Create document data for storage
                document_data = {
                    "name": filename,
                    "content": extracted_text,  # Store the extracted text
                    "mime_type": get_mime_type(file_ext),
                    "file_size": len(extracted_text.encode('utf-8'))  # Size of extracted text
                }
                
                # Save document using the provided db session
                save_document(chat_id, document_data, db)
                logger.info(f"Saved extracted text from '{filename}' to document store for chat {chat_id}")
            except Exception as save_error:
                logger.error(f"Failed to save document to store: {save_error}", exc_info=True)
                # Continue processing, don't fail just because storage failed

    except Exception as e:
        logger.error(f"Unexpected error processing file {filename}: {e}", exc_info=True)
        # Return empty string or re-raise depending on desired behavior
        extracted_text = f"[Unexpected error processing file: {str(e)}]"

    return extracted_text

def process_file_context(files_data: Optional[List[Dict[str, Any]]], chat_id: Optional[str] = None, db: Optional[Session] = None) -> str:
    """
    Processes a list of file dictionaries (from JSON payload) and concatenates their content.

    Args:
        files_data: A list of dictionaries, each expected to have 'name' and 'content' keys.
                    'content' should be string (text or base64 encoded for binary).
        chat_id: Optional chat ID to associate documents with for storage.
        db: Optional database session for storing documents.

    Returns:
        A single string containing the concatenated context from all processable files.
    """
    all_files_context = []
    if not files_data:
        logger.info("No files data provided for context processing.")
        return ""

    for file_dict in files_data:
        if isinstance(file_dict, dict) and file_dict.get("name"):
            file_context = process_file_content(file_dict, chat_id=chat_id, db=db)
            if file_context:
                # Use the filename from the dictionary
                filename = file_dict.get("name", "unknown_file")
                all_files_context.append(f"--- Start of File: {filename} ---\\n{file_context}\\n--- End of File: {filename} ---")
        else:
            logger.warning(f"Received invalid item in files data list: {type(file_dict)}. Skipping.")


    if not all_files_context:
        logger.info("No processable content found in the provided files.")
        return ""

    logger.info(f"Successfully processed content from {len(all_files_context)} file(s).")
    return "\\n\\n".join(all_files_context)

def save_document(chat_id: str, file_data: Dict[str, Any], db: Session) -> str:
    """
    Saves a document to the database, associating it with a chat ID.
    
    Args:
        chat_id: The ID of the chat to associate the document with
        file_data: Dictionary containing 'name' and 'content' keys
        db: Database session
    
    Returns:
        The ID of the created document
    """
    try:
        filename = file_data.get("name", "unknown_file")
        content_input = file_data.get("content")
        
        if not chat_id:
            raise ValueError("Chat ID is required")
        
        if not content_input:
            raise ValueError(f"No content provided for file: {filename}")
        
        # Generate a unique document ID
        document_id = str(uuid.uuid4())
        
        # Determine mime_type based on file extension
        file_ext = '.' + filename.split('.')[-1].lower() if '.' in filename else ''
        mime_type = get_mime_type(file_ext)
        
        # Calculate file size
        file_size = 0
        if isinstance(content_input, str):
            if file_ext in TEXT_EXTENSIONS:
                file_size = len(content_input.encode('utf-8'))
            else:
                # For base64 encoded files
                try:
                    file_bytes = base64.b64decode(content_input)
                    file_size = len(file_bytes)
                except Exception as e:
                    logger.error(f"Error calculating file size for {filename}: {e}")
                    file_size = len(content_input)
        
        # Insert into database using the provided session
        db.execute(text("""
            INSERT INTO chat.documents 
            (document_id, chat_id, filename, content, mime_type, file_size, created_at, updated_at)
            VALUES (:document_id, :chat_id, :filename, :content, :mime_type, :file_size, :created_at, :updated_at)
            RETURNING document_id
        """), {
            "document_id": document_id, 
            "chat_id": chat_id, 
            "filename": filename, 
            "content": content_input, 
            "mime_type": mime_type, 
            "file_size": file_size, 
            "created_at": datetime.now(), 
            "updated_at": datetime.now()
        })
        
        # Commit the changes
        db.commit()
        
        logger.info(f"Successfully saved document {filename} with ID {document_id} for chat {chat_id}")
        return document_id
        
    except Exception as e:
        db.rollback()
        logger.error(f"Error saving document for chat {chat_id}: {str(e)}", exc_info=True)
        raise

def get_document(document_id: str, db: Session) -> Dict[str, Any]:
    """
    Retrieves a document from the database by its ID.
    
    Args:
        document_id: The ID of the document to retrieve
        db: Database session
    
    Returns:
        Dictionary containing the document details
    """
    try:
        result = db.execute(text("""
            SELECT document_id, chat_id, filename, content, mime_type, file_size, created_at
            FROM chat.documents
            WHERE document_id = :document_id
        """), {"document_id": document_id}).fetchone()
        
        if not result:
            raise ValueError(f"Document with ID {document_id} not found")
        
        return {
            "document_id": result.document_id,
            "chat_id": result.chat_id,
            "filename": result.filename,
            "content": result.content,
            "mime_type": result.mime_type,
            "file_size": result.file_size,
            "created_at": result.created_at,
        }
        
    except Exception as e:
        logger.error(f"Error retrieving document {document_id}: {str(e)}", exc_info=True)
        raise

def get_chat_documents(chat_id: str, db: Session) -> List[Dict[str, Any]]:
    """
    Retrieves all documents associated with a chat.
    
    Args:
        chat_id: The ID of the chat
        db: Database session
    
    Returns:
        List of dictionaries containing document details (excluding content)
    """
    try:
        results = db.execute(text("""
            SELECT document_id, chat_id, filename, mime_type, file_size, created_at
            FROM chat.documents
            WHERE chat_id = :chat_id
            ORDER BY created_at DESC
        """), {"chat_id": chat_id}).fetchall()
        
        documents = []
        for result in results:
            documents.append({
                "document_id": result.document_id,
                "chat_id": result.chat_id,
                "filename": result.filename,
                "mime_type": result.mime_type,
                "file_size": result.file_size,
                "created_at": result.created_at,
            })
        
        return documents
        
    except Exception as e:
        logger.error(f"Error retrieving documents for chat {chat_id}: {str(e)}", exc_info=True)
        raise

def delete_document(document_id: str, db: Session) -> bool:
    """
    Deletes a document from the database.
    
    Args:
        document_id: The ID of the document to delete
        db: Database session
    
    Returns:
        True if successful, False otherwise
    """
    try:
        result = db.execute(text("""
            DELETE FROM chat.documents
            WHERE document_id = :document_id
            RETURNING document_id
        """), {"document_id": document_id}).fetchone()
        
        # Commit the changes
        db.commit()
        
        success = result is not None
        if success:
            logger.info(f"Successfully deleted document {document_id}")
        else:
            logger.warning(f"Document {document_id} not found for deletion")
        
        return success
        
    except Exception as e:
        db.rollback()
        logger.error(f"Error deleting document {document_id}: {str(e)}", exc_info=True)
        return False

def get_mime_type(file_ext: str) -> str:
    """
    Determines the MIME type based on file extension.
    
    Args:
        file_ext: File extension including the dot (e.g., '.pdf')
    
    Returns:
        The MIME type as a string
    """
    mime_types = {
        '.txt': 'text/plain',
        '.md': 'text/markdown',
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.rtf': 'application/rtf',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.csv': 'text/csv',
        '.png': 'image/png',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.gif': 'image/gif',
        '.svg': 'image/svg+xml',
    }
    
    return mime_types.get(file_ext.lower(), 'application/octet-stream') 